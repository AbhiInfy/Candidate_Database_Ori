from pathlib import Path
import subprocess

import imageio_ffmpeg
from PIL import Image, ImageDraw, ImageFilter, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT_DIR = ROOT / "Hiring Automation Tool"
ASSET_DIR = OUT_DIR / "client_ppt_assets"
OUT_PPT = OUT_DIR / "Motifzone Hiring Automation Tool - Client Presentation.pptx"

VIDEO = Path(r"C:\Users\Nitin Maheshwari\Videos\2026-06-01 10-38-00.mp4")
CAPTURE1 = Path(r"C:\Users\Nitin Maheshwari\Pictures\Capture1.PNG")
CAPTURE2 = Path(r"C:\Users\Nitin Maheshwari\Pictures\Capture2.PNG")
CAPTURE_UI = Path(r"C:\Users\Nitin Maheshwari\Pictures\Capture.PNG")
LOGO = Path(r"C:\Users\Nitin Maheshwari\Downloads\Logo\Logo\motifzone_020919_OP-01.jpg")
LOGO_ASSET = ASSET_DIR / "motifzone_logo_clean.png"

NAVY = RGBColor(16, 58, 81)
BRAND_BLUE = RGBColor(39, 47, 120)
BRAND_RED = RGBColor(239, 49, 37)
TEAL = RGBColor(36, 161, 156)
BG = RGBColor(246, 248, 251)
TEXT = RGBColor(20, 43, 60)
MUTED = RGBColor(84, 104, 120)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(42, 150, 92)
ORANGE = RGBColor(235, 137, 34)


def ffmpeg_extract(time_seconds: int, output: Path) -> None:
    ffmpeg = imageio_ffmpeg.get_ffmpeg_exe()
    subprocess.run(
        [
            ffmpeg,
            "-y",
            "-ss",
            str(time_seconds),
            "-i",
            str(VIDEO),
            "-frames:v",
            "1",
            str(output),
        ],
        check=True,
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
    )


def font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    candidates = [
        r"C:\Windows\Fonts\arialbd.ttf" if bold else r"C:\Windows\Fonts\arial.ttf",
        r"C:\Windows\Fonts\calibrib.ttf" if bold else r"C:\Windows\Fonts\calibri.ttf",
    ]
    for candidate in candidates:
        if Path(candidate).exists():
            return ImageFont.truetype(candidate, size)
    return ImageFont.load_default()


def redact_rect(image: Image.Image, box: tuple[int, int, int, int], label: str | None = None) -> None:
    crop = image.crop(box).filter(ImageFilter.GaussianBlur(9))
    image.paste(crop, box)
    draw = ImageDraw.Draw(image)
    draw.rounded_rectangle(box, radius=4, fill=(235, 242, 247), outline=(199, 214, 225), width=2)
    if label:
        draw.text((box[0] + 8, box[1] + 8), label, fill=(16, 58, 81), font=font(18, True))


def prepare_logo() -> Path:
    image = Image.open(LOGO).convert("RGB")
    # Crop large white margins from the source logo while preserving the white logo background.
    gray = image.convert("L")
    mask = gray.point(lambda p: 0 if p > 246 else 255)
    bbox = mask.getbbox()
    if bbox:
        pad = 60
        left = max(bbox[0] - pad, 0)
        top = max(bbox[1] - pad, 0)
        right = min(bbox[2] + pad, image.width)
        bottom = min(bbox[3] + pad, image.height)
        image = image.crop((left, top, right, bottom))
    image.thumbnail((900, 360), Image.Resampling.LANCZOS)
    canvas = Image.new("RGB", (920, 380), "white")
    canvas.paste(image, ((canvas.width - image.width) // 2, (canvas.height - image.height) // 2))
    canvas.save(LOGO_ASSET)
    return LOGO_ASSET


def prepare_excel_screenshot() -> Path:
    image = Image.open(CAPTURE1).convert("RGB")
    # Mask candidate personal data while retaining the Excel output structure.
    redact_rect(image, (35, 75, 226, 615), "Candidate names masked")
    redact_rect(image, (232, 75, 384, 615), "Emails masked")
    redact_rect(image, (388, 75, 580, 615), "Phones masked")
    redact_rect(image, (1330, 75, 1715, 615), "CV file paths masked")
    draw = ImageDraw.Draw(image)
    draw.rounded_rectangle((1260, 25, 1845, 70), radius=8, fill=(16, 58, 81))
    draw.text((1280, 36), "Excel tracker: structured candidate data", fill="white", font=font(24, True))
    path = ASSET_DIR / "masked_excel_output.png"
    image.save(path)
    return path


def prepare_download_screenshot() -> Path:
    image = Image.open(CAPTURE2).convert("RGB")
    redact_rect(image, (95, 64, 485, 254), "CV filenames masked")
    draw = ImageDraw.Draw(image)
    draw.rounded_rectangle((520, 54, 1080, 105), radius=8, fill=(16, 58, 81))
    draw.text((545, 68), "Downloaded CV folder: PDF and Word resumes", fill="white", font=font(24, True))
    path = ASSET_DIR / "masked_downloaded_cvs.png"
    image.save(path)
    return path


def prepare_ui_screenshot() -> Path:
    image = Image.open(CAPTURE_UI).convert("RGB")
    # Hide local path details in the run log while retaining the successful run summary.
    draw = ImageDraw.Draw(image)
    redact_rect(image, (360, 842, 1140, 898), "Local file paths masked")
    draw.rounded_rectangle((1020, 55, 1438, 110), radius=8, fill=(16, 58, 81))
    draw.text((1042, 70), "Live run completed successfully", fill="white", font=font(24, True))
    path = ASSET_DIR / "live_run_success.png"
    image.save(path)
    return path


def make_flow_image() -> Path:
    path = ASSET_DIR / "workflow.png"
    image = Image.new("RGB", (1600, 680), "#f6f8fb")
    draw = ImageDraw.Draw(image)
    steps = [
        ("1. Run Tool", "User enters app\npassword and clicks run", 70, 185),
        ("2. Scan Mailbox", "Reads new emails\nfrom hiring inbox", 525, 185),
        ("3. Filter CVs", "Downloads only\nresume attachments", 980, 185),
        ("4. Extract Data", "Name, email, phone,\nexperience and skills", 980, 445),
        ("5. Check Duplicates", "Compares existing\ncandidate records", 525, 445),
        ("6. Update Excel", "Formatted tracker\nready for review", 70, 445),
    ]
    box_w = 330
    box_h = 140
    draw.text((70, 42), "Automation Workflow", fill="#103a51", font=font(52, True))
    draw.text((74, 108), "Clear end-to-end flow from mailbox to Excel tracker", fill="#355266", font=font(27))

    def arrow(start: tuple[int, int], end: tuple[int, int]) -> None:
        sx, sy = start
        ex, ey = end
        draw.line((sx, sy, ex, ey), fill="#24a19c", width=9)
        if abs(ex - sx) > abs(ey - sy):
            direction = 1 if ex > sx else -1
            draw.polygon([(ex, ey), (ex - 28 * direction, ey - 18), (ex - 28 * direction, ey + 18)], fill="#24a19c")
        else:
            direction = 1 if ey > sy else -1
            draw.polygon([(ex, ey), (ex - 18, ey - 28 * direction), (ex + 18, ey - 28 * direction)], fill="#24a19c")

    connectors = [
        ((400, 255), (510, 255)),
        ((855, 255), (965, 255)),
        ((1145, 325), (1145, 430)),
        ((980, 515), (870, 515)),
        ((525, 515), (415, 515)),
    ]
    for start, end in connectors:
        arrow(start, end)

    for idx, (title, body, x, y) in enumerate(steps):
        fill = "#103a51" if idx in (0, 5) else "#ffffff"
        outline = "#ef3125" if idx in (0, 5) else "#24a19c"
        title_fill = "#ffffff" if idx in (0, 5) else "#103a51"
        body_fill = "#d8f2ef" if idx in (0, 5) else "#355266"
        draw.rounded_rectangle((x, y, x + box_w, y + box_h), radius=18, fill=fill, outline=outline, width=4)
        draw.text((x + 24, y + 22), title, fill=title_fill, font=font(28, True))
        draw.multiline_text((x + 24, y + 68), body, fill=body_fill, font=font(22), spacing=6)
    image.save(path)
    return path


def prepare_assets() -> dict[str, Path]:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    assets = {
        "logo": prepare_logo(),
        "excel": prepare_excel_screenshot(),
        "downloads": prepare_download_screenshot(),
        "ui_capture": prepare_ui_screenshot(),
        "flow": make_flow_image(),
    }
    for seconds, name in [(4, "video_start"), (15, "video_running"), (30, "video_result")]:
        out = ASSET_DIR / f"{name}.png"
        ffmpeg_extract(seconds, out)
        assets[name] = out
    return assets


def add_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_header(slide, title: str, subtitle: str | None = None) -> None:
    add_bg(slide)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.95))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    tf = slide.shapes.add_textbox(Inches(0.55), Inches(0.2), Inches(7.2), Inches(0.36)).text_frame
    tf.text = "Motifzone Private Limited"
    tf.paragraphs[0].runs[0].font.size = Pt(18)
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.color.rgb = WHITE
    if LOGO_ASSET.exists():
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(11.05), Inches(0.17), Inches(1.65), Inches(0.56))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.fill.background()
        slide.shapes.add_picture(str(LOGO_ASSET), Inches(11.16), Inches(0.27), width=Inches(1.42))
    tf2 = slide.shapes.add_textbox(Inches(0.55), Inches(1.18), Inches(9.9), Inches(0.48)).text_frame
    tf2.text = title
    tf2.paragraphs[0].runs[0].font.size = Pt(28)
    tf2.paragraphs[0].runs[0].font.bold = True
    tf2.paragraphs[0].runs[0].font.color.rgb = NAVY
    if subtitle:
        st = slide.shapes.add_textbox(Inches(0.58), Inches(1.72), Inches(11.8), Inches(0.35)).text_frame
        st.text = subtitle
        st.paragraphs[0].runs[0].font.size = Pt(15)
        st.paragraphs[0].runs[0].font.color.rgb = MUTED


def set_text(shape, text: str, size: int = 18, color=TEXT, bold: bool = False, align=None) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    if align:
        p.alignment = align


def add_pill(slide, x, y, w, text, color=TEAL) -> None:
    pill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.38))
    pill.fill.solid()
    pill.fill.fore_color.rgb = color
    pill.line.fill.background()
    set_text(pill, text, 11, WHITE, True, PP_ALIGN.CENTER)
    pill.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_metric(slide, x, y, number, label, color=TEAL) -> None:
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.45), Inches(1.35))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = RGBColor(216, 226, 235)
    set_text(box, f"{number}\n{label}", 18, color, True, PP_ALIGN.CENTER)
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_bullet_box(slide, x, y, w, h, title, bullets, accent=TEAL) -> None:
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = RGBColor(216, 226, 235)
    title_box = slide.shapes.add_textbox(Inches(x + 0.25), Inches(y + 0.2), Inches(w - 0.5), Inches(0.33))
    set_text(title_box, title, 16, NAVY, True)
    yy = y + 0.66
    for bullet in bullets:
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.28), Inches(yy + 0.08), Inches(0.11), Inches(0.11))
        dot.fill.solid()
        dot.fill.fore_color.rgb = accent
        dot.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(x + 0.48), Inches(yy - 0.02), Inches(w - 0.75), Inches(0.5))
        set_text(tb, bullet, 12, TEXT)
        yy += 0.53


def add_image(slide, path: Path, x, y, w, h=None) -> None:
    pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h) if h else None)
    pic.line.color.rgb = RGBColor(199, 214, 225)
    return pic


def add_logo_card(slide, x, y, w) -> None:
    if not LOGO_ASSET.exists():
        return
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(w * 0.42))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(235, 240, 245)
    slide.shapes.add_picture(str(LOGO_ASSET), Inches(x + 0.12), Inches(y + 0.07), width=Inches(w - 0.24))


def build_deck() -> None:
    assets = prepare_assets()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1. Title
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    block = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    block.fill.solid()
    block.fill.fore_color.rgb = NAVY
    block.line.fill.background()
    add_logo_card(slide, 0.75, 0.72, 2.35)
    title = slide.shapes.add_textbox(Inches(0.78), Inches(1.95), Inches(6.0), Inches(1.0))
    set_text(title, "Hiring Automation Tool", 42, WHITE, True)
    subtitle = slide.shapes.add_textbox(Inches(0.82), Inches(3.0), Inches(5.9), Inches(0.78))
    set_text(subtitle, "Automated resume email processing and candidate Excel tracking", 20, RGBColor(214, 239, 240))
    add_pill(slide, 0.84, 4.08, 2.05, "Client Demo", BRAND_RED)
    add_pill(slide, 3.05, 4.08, 2.35, "Motifzone Pvt Ltd", BRAND_BLUE)
    add_image(slide, assets["ui_capture"], 7.1, 0.95, 5.45, 3.25)
    add_image(slide, assets["excel"], 7.1, 4.48, 5.45, 1.85)

    # 2. Requirement
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Business Requirement", "Recruiters need faster, cleaner handling of CVs received by email.")
    add_bullet_box(
        slide,
        0.7,
        2.35,
        5.8,
        3.35,
        "Current Manual Pain Points",
        [
            "Download every attachment manually from the mailbox.",
            "Open each CV and copy candidate details one by one.",
            "Risk of duplicate entries and missed resumes.",
            "Excel formatting and sorting takes extra time.",
        ],
        ORANGE,
    )
    add_bullet_box(
        slide,
        6.9,
        2.35,
        5.7,
        3.35,
        "Automation Goal",
        [
            "Read incoming CV emails from the hiring mailbox.",
            "Download only valid resume attachments.",
            "Extract important candidate information.",
            "Maintain one clean Excel tracker for the recruitment team.",
        ],
        TEAL,
    )

    # 3. Solution overview
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Solution Overview", "A user-friendly desktop tool hides the Python code and runs the automation from one button.")
    add_image(slide, assets["video_start"], 0.72, 2.22, 6.1)
    add_bullet_box(
        slide,
        7.25,
        2.18,
        5.35,
        3.55,
        "What the Tool Does",
        [
            "Connects to careers@emotifzone.com through IMAP.",
            "Processes new emails incrementally after the first setup.",
            "Shows skipped JD/non-CV files clearly for client visibility.",
            "Keeps candidate output in a formatted Excel file.",
        ],
    )

    # 4. Workflow
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Process Flow", "From mailbox attachment to recruiter-ready candidate tracker.")
    add_image(slide, assets["flow"], 0.62, 2.05, 12.05, 5.1)

    # 5. Live run
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Live Run Screen", "The user sees a simple interface while the automation runs in the background.")
    add_image(slide, assets["ui_capture"], 0.7, 2.02, 7.25, 4.85)
    add_bullet_box(
        slide,
        8.25,
        2.2,
        4.35,
        3.35,
        "Visible Progress",
        [
            "Run New Emails for incremental processing.",
            "Full Scan option when historical mailbox scan is needed.",
            "Skipped JD attachments are shown in the log.",
            "Success message confirms Excel and CV folder output.",
        ],
    )

    # 6. Downloaded CVs
    slide = prs.slides.add_slide(blank)
    add_header(slide, "CV Download Automation", "The tool downloads only valid CV files into a clean folder.")
    add_image(slide, assets["downloads"], 0.72, 2.18, 7.65)
    add_bullet_box(
        slide,
        8.65,
        2.25,
        3.95,
        3.25,
        "Supported Files",
        [
            "PDF resumes",
            "Microsoft Word resumes",
            "Scanned PDFs with OCR support when available",
            "JD/non-CV files are skipped",
        ],
    )

    # 7. Excel output
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Excel Candidate Tracker", "Extracted resume data is stored in a structured Excel sheet.")
    add_image(slide, assets["excel"], 0.58, 2.02, 12.15)
    add_pill(slide, 0.82, 6.55, 2.15, "Candidate Name", BRAND_BLUE)
    add_pill(slide, 3.2, 6.55, 1.5, "Email", BRAND_BLUE)
    add_pill(slide, 4.95, 6.55, 2.0, "Contact No.", BRAND_BLUE)
    add_pill(slide, 7.2, 6.55, 1.55, "Skills", BRAND_BLUE)
    add_pill(slide, 9.0, 6.55, 2.15, "CV File Link", BRAND_BLUE)

    # 8. Controls and data quality
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Controls And Data Quality", "Built-in checks keep the tracker cleaner and easier to review.")
    add_metric(slide, 0.9, 2.35, "No", "Duplicate Rows", TEAL)
    add_metric(slide, 3.7, 2.35, "New", "Emails Only", GREEN)
    add_metric(slide, 6.5, 2.35, "PDF + DOC", "CV Support", NAVY)
    add_metric(slide, 9.3, 2.35, "Excel", "Formatted Output", ORANGE)
    add_bullet_box(
        slide,
        1.0,
        4.35,
        11.35,
        1.55,
        "Data Handling",
        [
            "Duplicate checking uses candidate name, email and phone where available.",
            "Excel headings are formatted, filters are available, and records are sorted by latest received date.",
        ],
    )

    # 9. Benefits
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Business Benefits", "A practical achievement that reduces repetitive recruiter effort.")
    add_bullet_box(
        slide,
        0.8,
        2.15,
        3.75,
        3.55,
        "Time Saving",
        [
            "Bulk CV download is automatic.",
            "Candidate data extraction is faster.",
            "Excel tracking is ready after each run.",
        ],
        GREEN,
    )
    add_bullet_box(
        slide,
        4.85,
        2.15,
        3.75,
        3.55,
        "Process Quality",
        [
            "Reduces copy-paste errors.",
            "Avoids repeated candidate rows.",
            "Separates CVs from JD attachments.",
        ],
        TEAL,
    )
    add_bullet_box(
        slide,
        8.9,
        2.15,
        3.75,
        3.55,
        "Client Value",
        [
            "Simple desktop experience.",
            "Clear run status and output.",
            "Professional recruitment tracker.",
        ],
        ORANGE,
    )

    # 10. Closing
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    add_logo_card(slide, 0.75, 0.72, 2.35)
    title = slide.shapes.add_textbox(Inches(0.78), Inches(2.0), Inches(6.0), Inches(0.8))
    set_text(title, "Hiring Automation Tool", 42, WHITE, True)
    sub = slide.shapes.add_textbox(Inches(0.82), Inches(2.9), Inches(5.8), Inches(0.8))
    set_text(sub, "A ready-to-use automation for resume mailbox processing and candidate tracking.", 20, RGBColor(214, 239, 240))
    add_image(slide, assets["video_result"], 7.1, 0.95, 5.45)
    add_bullet_box(
        slide,
        0.85,
        4.05,
        5.65,
        2.1,
        "Recommended Next Steps",
        [
            "Run the tool daily for new candidate emails.",
            "Review the Excel tracker before sharing externally.",
            "Continue improving name extraction with real CV examples.",
        ],
    )
    footer = slide.shapes.add_textbox(Inches(0.85), Inches(6.72), Inches(6.8), Inches(0.3))
    set_text(footer, "Prepared for client demonstration by Motifzone Private Limited", 13, RGBColor(214, 239, 240))

    prs.save(OUT_PPT)
    print(OUT_PPT)


if __name__ == "__main__":
    build_deck()
